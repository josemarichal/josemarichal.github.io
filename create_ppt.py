from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Define a simple slide layout function
    def add_slide(title_text, content_text):
        slide_layout = prs.slide_layouts[1] # Content with Caption layout or similar
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        
        # Body
        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.text = content_text
        
        # Formatting adjustments could be added here
        return slide

    # Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Politics in the Algorithmic Age"
    subtitle = slide.placeholders[1]
    subtitle.text = "Course: Global Politics 101\nTopic: The Algorithmic Social Contract"

    # Slide 1
    add_slide(
        "The Algorithmic Age",
        "Entering the Algorithmic Age\n\n"
        "• Definition: Digital interactions (likes, swipes) broken into Nodes (things) and Edges (attributes).\n"
        "• The Shift: Pre-2010 focused on understanding humans. Post-2010 focuses on predicting behavior.\n"
        "• Key Takeaway: Transition from surveillance for marketing to predictive modeling for behavioral replication."
    )

    # Slide 2
    add_slide(
        "The Logic of the Social Contract",
        "Political Legitimacy and Authority\n\n"
        "• Theory: Rational actors surrender freedom to authority for benefits.\n"
        "• Hobbes: Cede rights for protection (security).\n"
        "• Locke: Cede power to preserve rights (life, liberty, property).\n"
        "• Rousseau: Submit to 'General Will' for meaning.\n"
        "• Relevance: Explains our relationship with digital platforms."
    )

    # Slide 3
    add_slide(
        "The Algorithmic Social Contract",
        "Terms of the New Deal\n\n"
        "• The Problem: Infinite information causes 'Anxiety of Choice'.\n"
        "• The Exchange: We cede curational autonomy to algorithms.\n"
        "• The Benefit: Reduced anxiety, comfort, filtered dissonance.\n"
        "• The Cost: Data extraction, surveillance, consumer clustering."
    )

    # Slide 4
    add_slide(
        "The Paradox of Expression",
        "Voice vs. Reflection\n\n"
        "• The Promise: Early internet promised authentic self-expression.\n"
        "• The Reality: Must modify voice to be 'heard' by algorithms.\n"
        "• Voice (Monetizable): Immediate, instinctive, reactionary.\n"
        "• Reflection (Hard to Monetize): Slow, thoughtful.\n"
        "• Result: Malleable identity performed for engagement."
    )

    # Slide 5
    add_slide(
        "The Economics of Engagement",
        "Opinion as Supply and Demand\n\n"
        "• Datafication: Human expression tokenized for AI training.\n"
        "• Market Incentive: Constant supply of 'voice' needed.\n"
        "• Habitual Engagement: Production of opinion becomes habit.\n"
        "• The Intelligence Gap: Bots vs Humans matters less than the data generated."
    )

    # Slide 6
    add_slide(
        "The Outlier Problem",
        "Map vs. Territory\n\n"
        "• Traditional Science: Simple models (parsimony), accepted outliers.\n"
        "• AI Revolution: Billions of parameters to force 'fit'.\n"
        "• Goal: Prediction (what happens next) over Explanation (why).\n"
        "• The Friction: Humans are complex/contingent and don't fit static models."
    )

    # Slide 7
    add_slide(
        "Machine Habitus",
        "Conforming to the Code\n\n"
        "• The Danger: If we are too unpredictable, we must change.\n"
        "• Machine Habitus: Subconsciously sorting/classifying ourselves to match the algorithm.\n"
        "• Optimization: Moving towards the 'Local Minima' - adapting to the model's simplicity."
    )

    prs.save('Algorithmic_Politics_Slides.pptx')
    print("Successfully created 'Algorithmic_Politics_Slides.pptx'")

if __name__ == "__main__":
    create_presentation()
