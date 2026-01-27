from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_presentation():
    prs = Presentation()

    # Define a helper function to add slides with images and provocations
    def add_slide_with_layout(title_text, bullet_points, image_filename, provocation_text):
        # Using a blank layout to have full control
        slide = prs.slides.add_slide(prs.slide_layouts[6]) 

        # Add Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        title_tf = title_box.text_frame
        title_tf.text = title_text
        title_p = title_tf.paragraphs[0]
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.name = 'Arial'

        # Add Bullet Points (Left Side)
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(4.5))
        tf = text_box.text_frame
        tf.word_wrap = True
        
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(16)
            p.space_after = Pt(10)
            p.level = 0
            # Basic bullet formatting happens automatically or can be customized further

        # Add Image (Right Side)
        if image_filename:
            image_path = os.path.join("images", image_filename)
            if os.path.exists(image_path):
                # Add picture, resizing width to fit nicely on the right
                slide.shapes.add_picture(image_path, Inches(5.2), Inches(1.5), width=Inches(4.5))

        # Add Provocation (Bottom)
        if provocation_text:
            prov_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
            prov_tf = prov_box.text_frame
            p = prov_tf.paragraphs[0]
            p.text = provocation_text
            p.font.size = Pt(20)
            p.font.italic = True
            p.font.color.rgb = RGBColor(0, 112, 192) # A blue-ish color
            p.alignment = PP_ALIGN.CENTER

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Politics in the Algorithmic Age"
    subtitle = slide.placeholders[1]
    subtitle.text = "Course: Global Politics 101\nTopic: The Algorithmic Social Contract"

    # Slide 1
    add_slide_with_layout(
        "The Algorithmic Age",
        [
            "Definition: We live in an 'Algorithmic Age' where every digital interaction is broken down into Nodes (things) and Edges (attributes).",
            "The Shift: Prior to 2010, the goal was to understand humans. Now, the priority is to predict future behavior.",
            "Key Takeaway: Transition from surveillance for marketing to predictive modeling for behavioral replication."
        ],
        "slide_1.png",
        "\"If you can be predicted, are you free?\""
    )

    # Slide 2
    add_slide_with_layout(
        "The Logic of the Social Contract",
        [
            "Theory: Rational actors surrender freedom to a political authority for benefits.",
            "Hobbes: We cede rights for protection (security).",
            "Locke: We cede power to preserve rights (life, liberty).",
            "Rousseau: We submit to the 'General Will' for meaning."
        ],
        "slide_2.png",
        "\"What rights did you sign away today?\""
    )

    # Slide 3
    add_slide_with_layout(
        "The Algorithmic Social Contract",
        [
            "The Problem: The internet offers infinite information ('Anxiety of Choice').",
            "The Exchange: We cede curational autonomy to algorithms.",
            "The Benefit: Curated comfort and reduced anxiety.",
            "The Cost: Data extraction and consumer clustering."
        ],
        "slide_3.png",
        "\"Is comfort worth the cost of your autonomy?\""
    )

    # Slide 4
    add_slide_with_layout(
        "The Paradox of Expression",
        [
            "The Promise: The internet as a place for the 'authentic self.'",
            "The Reality: Modifying our voice to be 'heard' by the code.",
            "Reflection vs. Voice: Reflection is slow and hard to monetize. Voice is immediate and profitable."
        ],
        "slide_4.png",
        "\"Are you sharing your life, or performing it?\""
    )

    # Slide 5
    add_slide_with_layout(
        "The Economics of Engagement",
        [
            "Datafication: Human expression tokenized for AI training.",
            "Supply & Demand: Platforms need a constant supply of opinion.",
            "Habitual Engagement: 'Voice' becomes a habit, not a thought."
        ],
        "slide_5.png",
        "\"Who profits from your outrage?\""
    )

    # Slide 6
    add_slide_with_layout(
        "The Outlier Problem",
        [
            "Traditional Science: Accepted 'outliers' (parsimony).",
            "AI Revolution: Uses billions of parameters to force the model to fit reality.",
            "Goal: Prediction (what happens next) over Explanation (why)."
        ],
        "slide_6.png",
        "\"What parts of you don't fit the model?\""
    )

    # Slide 7
    add_slide_with_layout(
        "Machine Habitus",
        [
            "The Danger: If we are unpredictable, the system is inefficient.",
            "Machine Habitus: We subconsciously sort and classify ourselves to match the algorithm.",
            "Conclusion: We adapt to the machine's simplicity instead of it adapting to our complexity."
        ],
        "slide_7.png",
        "\"Are you training the algorithm, or is it training you?\""
    )

    output_file = 'Algorithmic_Politics_Full.pptx'
    prs.save(output_file)
    print(f"Successfully created '{output_file}'")

if __name__ == "__main__":
    create_presentation()
